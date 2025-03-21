import os
import fitz  # PyMuPDF
import docx
import pptx
import requests
from typing import TypedDict

from langgraph.graph import StateGraph
from langchain_community.chat_models import ChatOpenAI
from langchain.schema import HumanMessage, SystemMessage
from langchain_core.runnables import RunnableLambda

from bs4 import BeautifulSoup

def fetch_text_from_links(links):
    texts = []
    MAX_CHARS = 8000

    for url in links:
        try:
            response = requests.get(url, timeout=10)
            soup = BeautifulSoup(response.text, 'html.parser')

            main_content = soup.find('div', {'id': 'bodyContent'})
            clean_text = main_content.get_text(separator="\n") if main_content else soup.get_text()

            if len(clean_text) > MAX_CHARS:
                print(f"⚠️ Trimming content from {len(clean_text)} to {MAX_CHARS} characters.")
                clean_text = clean_text[:MAX_CHARS]

            print(f"📝 Raw text from {url}:\n{clean_text[:500]}...\n---")
            texts.append(clean_text)
        except Exception as e:
            print(f"⚠️ Failed to fetch {url}: {e}")
    return "\n".join(texts)


def summarize_urls(request, links):
    url_summaries = {}
    combined_llm = get_llm(request)
    for url in links:
        try:
            raw_text = fetch_text_from_links([url])
            chunks = extract_text_chunks(raw_text)
            summaries = []
            for chunk in chunks:
                resp = combined_llm.invoke([HumanMessage(content=f"Context:\n{chunk}\n\nSummarize this for knowledge extraction.")])
                summaries.append(resp.content)
            url_summaries[url] = summaries
        except Exception as e:
            print(f"❌ Failed to summarize {url}: {e}")
    request.session["url_summaries"] = url_summaries
    request.session.modified = True


# ---------------------- TEXT EXTRACTION ----------------------
def extract_text_from_file(file_path):
    if file_path.endswith(".txt"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    elif file_path.endswith(".pdf"):
        text = ""
        with fitz.open(file_path) as pdf:
            for page in pdf:
                text += page.get_text()
        return text

    elif file_path.endswith(".docx"):
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    elif file_path.endswith(".pptx"):
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    return "Unsupported file type."


def extract_text_chunks(text, chunk_size=5000):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]


# ---------------------- LLM WRAPPERS ----------------------
class MistralWrapper:
    def __init__(self, api_key, system_prompt=None):
        self.api_key = api_key
        self.system_prompt = system_prompt or ""
        self.endpoint = "https://api.mistral.ai/v1/chat/completions"

    def invoke(self, messages):
        content = "\n".join([msg.content for msg in messages if isinstance(msg, HumanMessage)])
        all_messages = []
        if self.system_prompt:
            all_messages.append({"role": "system", "content": self.system_prompt})
        all_messages.append({"role": "user", "content": content})

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json",
        }
        data = {
            "model": "mistral-large-latest",
            "messages": all_messages
        }
        response = requests.post(self.endpoint, headers=headers, json=data)
        response.raise_for_status()
        return type('LLMResponse', (object,), {"content": response.json()["choices"][0]["message"]["content"]})()


# ---------------------- LLM SETUP ----------------------
def get_llm(request_or_dict):
    if isinstance(request_or_dict, dict):
        session = request_or_dict.get("session", {})
    else:
        session = getattr(request_or_dict, 'session', {})

    model = session.get("model_name")
    api_key = session.get("api_key")
    rules = session.get("rules", [])

    if not api_key:
        raise ValueError("API key is required.")

    system_prompt = "\n".join(rules) if rules else None

    if model == "llama3":
        llm = ChatOpenAI(
            model_name="llama3-8b-8192",
            openai_api_key=api_key,
            base_url="https://api.groq.com/openai/v1"
        )
    elif model == "gpt4":
        llm = ChatOpenAI(
            model_name="gpt-4",
            openai_api_key=api_key,
        )
    elif model == "mistral":
        return MistralWrapper(api_key, system_prompt)
    else:
        raise ValueError("Invalid model selected.")

    def run_with_system(messages):
        if system_prompt:
            messages = [SystemMessage(content=system_prompt)] + messages
        return llm.invoke(messages)

    return RunnableLambda(run_with_system)


# ---------------------- AGENT FUNCTIONS ----------------------
def retrieve_data(state):
    request = state["request"]
    llm = get_llm(request)
    query = state["input_text"]
    file_path = state.get("file_path")

    file_text = request.session.get("file_text")
    file_summaries = request.session.get("file_summaries")

    # ✅ Use precomputed summaries if available
    if file_summaries:
        print("⚡ Using precomputed summaries from session")
        return {"retrieved_data": "\n---\n".join(file_summaries)}

    # 🔁 Fallback to processing chunks
    if not file_text and file_path and os.path.exists(file_path):
        print("📂 Re-reading file from disk")
        file_text = extract_text_from_file(file_path)
        request.session["file_text"] = file_text
        request.session.modified = True

    if file_text:
        print("📄 Processing file content with chunking")
        chunks = extract_text_chunks(file_text)
        summaries = []

        for i, chunk in enumerate(chunks):
            print(f"🔍 Processing chunk {i+1}/{len(chunks)}")
            resp = llm.invoke([
                HumanMessage(content=f"Context:\n{chunk}\n\nBased on this context, answer or extract information related to: {query}")
            ])
            summaries.append(resp.content)

        request.session["file_summaries"] = summaries
        request.session.modified = True

        return {"retrieved_data": "\n---\n".join(summaries)}

    print("⚠️ No file content available. Using pure LLM.")
    response = llm.invoke([HumanMessage(content=f"{query}")])
    return {"retrieved_data": response.content}


def summarize_data(state):
    request = state["request"]
    llm = get_llm(request)
    response = llm.invoke([HumanMessage(content=f"Summarize this: {state['retrieved_data']}")])
    return {"summarized_data": response.content}


def answer_question(state):
    request = state["request"]
    llm = get_llm(request)

    context_parts = []

    if "summarized_data" in state and state["summarized_data"]:
        context_parts.append("File Context:\n" + state["summarized_data"])

    if "url_context" in state and state["url_context"].strip():
        context_parts.append("URL Context:\n" + state["url_context"].strip())

    final_context = "\n\n".join(context_parts).strip()

    if not final_context:
        prompt = f"Answer the following question using only your general knowledge:\n\nQ: {state['input_text']}"
    else:
        prompt = f"Based on the following context, answer the question accurately:\n\n{final_context}\n\nQ: {state['input_text']}"

    response = llm.invoke([HumanMessage(content=prompt)])

    return {"final_answer": response.content}





# ---------------------- STATE GRAPH ----------------------
class AgentState(TypedDict):
    input_text: str
    file_path: str
    request: any
    retrieved_data: str
    summarized_data: str
    final_answer: str


def process_file_with_graph(request, input_text, file_path=None):
    # Prepare RAG from URLs
    url_summaries = request.session.get("url_summaries", {})
    url_context = ""

    if url_summaries:
        for url, summaries in url_summaries.items():
            url_context += f"\n[From URL: {url}]\n" + "\n".join(summaries) + "\n"

    # Build the workflow
    workflow = StateGraph(AgentState)

    workflow.add_node("retriever", retrieve_data)
    workflow.add_node("summarizer", summarize_data)
    workflow.add_node("qa_agent", answer_question)

    workflow.add_edge("retriever", "summarizer")
    workflow.add_edge("summarizer", "qa_agent")
    workflow.set_entry_point("retriever")

    executor = workflow.compile()

    # Inject context into the state
    input_state = {
        "input_text": input_text,
        "file_path": file_path or "",
        "request": request,
        "url_context": url_context.strip(),  # 🌐 added URL summaries here
    }

    result = executor.invoke(input_state)
    return result
